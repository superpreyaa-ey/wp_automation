from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.util import Inches, Pt
import pandas as pd
# Convert hex color to RGB
def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def create_presentation_report(file_path,filename, bg_color_hex):

    # Create a presentation object
    xls = pd.ExcelFile(file_path)
    audit_df = pd.read_excel(xls, sheet_name='Audit')
    issue_df = pd.read_excel(xls, sheet_name='Issues')
    
    prs = Presentation()

    # Define the background color
    bg_color = hex_to_rgb(bg_color_hex)

    # Add a slide with a blank layout (no placeholders)
    slide_layout = prs.slide_layouts[6]  # 6 is the layout for a blank slide
    slide = prs.slides.add_slide(slide_layout)

    # Access the background of the slide
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Add title text to the first slide with custom font size and style
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    text_frame = title_box.text_frame
    p = text_frame.add_paragraph()
    p.text = "Hello, World!"
    p.font.size = Pt(28)  # Set font size to 28 points
    p.font.name = 'PrudentialModern SemCond'  # Set font name to 'PrudentialModern SemCond'
    p.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

    # Apply the font style to the entire text in the text box
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(28)
        paragraph.font.name = 'PrudentialModern SemCond'
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Save the presentation
    prs.save(filename)

# Example usage
# file_path = 'C:/Prudential/Code/PROUD_Automation/Compliance/static/media/Report/Audit Report.xlsx'
# ret_val = create_presentation_report(file_path,'example_with_bg_and_style.pptx', '#1a1a24')
# print(ret_val)
# Save the presentation
# prs.save('example.pptx')
#1a1a24